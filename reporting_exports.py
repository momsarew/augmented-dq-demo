cat > reporting_exports.py << 'EOF'
from __future__ import annotations

from io import BytesIO
from datetime import datetime
from typing import Dict, Any, List

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

from openpyxl import Workbook
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches, Pt


def export_report_pdf(markdown_text: str, title: str = "Rapport Qualité Données") -> bytes:
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=A4)
    width, height = A4

    c.setTitle(title)
    c.setFont("Helvetica-Bold", 14)
    c.drawString(40, height - 50, title)
    c.setFont("Helvetica", 9)
    c.drawString(40, height - 65, datetime.now().strftime("Généré le %Y-%m-%d %H:%M"))

    y = height - 90
    c.setFont("Courier", 8)

    for line in markdown_text.splitlines():
        if y < 60:
            c.showPage()
            y = height - 50
            c.setFont("Courier", 8)

        if len(line) > 130:
            chunks = [line[i:i+130] for i in range(0, len(line), 130)]
            for ch in chunks:
                c.drawString(40, y, ch)
                y -= 10
        else:
            c.drawString(40, y, line)
            y -= 10

    c.save()
    buffer.seek(0)
    return buffer.getvalue()


def export_report_excel(
    results: Dict[str, Any],
    profil: str,
    attribut_focus: str,
    usage_focus: str
) -> bytes:
    wb = Workbook()

    ws = wb.active
    ws.title = "Meta"
    ws.append(["Profil", profil])
    ws.append(["Attribut focus", attribut_focus])
    ws.append(["Usage focus", usage_focus])
    ws.append(["Date", datetime.now().isoformat(timespec="minutes")])

    ws2 = wb.create_sheet("Vecteur_4D")
    ws2.append(["Attribut", "P_DB", "P_DP", "P_BR", "P_UP"])
    vect = results.get("vecteurs_4d", {}).get(attribut_focus, {}) or {}
    ws2.append([
        attribut_focus,
        float(vect.get("P_DB", 0)),
        float(vect.get("P_DP", 0)),
        float(vect.get("P_BR", 0)),
        float(vect.get("P_UP", 0)),
    ])

    ws3 = wb.create_sheet("Scores")
    ws3.append(["Attribut_Usage", "Score"])
    scores = results.get("scores", {}) or {}
    for k, v in sorted(scores.items(), key=lambda x: float(x[1] or 0), reverse=True)[:50]:
        ws3.append([k, float(v or 0)])

    ws4 = wb.create_sheet("Top_Priorites")
    ws4.append(["#", "Attribut", "Usage", "Score", "Severite", "Records_affected", "Impact_mensuel"])
    prios = results.get("top_priorities", []) or []
    for i, p in enumerate(prios[:20], 1):
        ws4.append([
            i,
            p.get("attribut", ""),
            p.get("usage", ""),
            float(p.get("score", 0) or 0),
            p.get("severite", ""),
            p.get("records_affected", ""),
            p.get("impact_mensuel", ""),
        ])

    for sheet in [ws, ws2, ws3, ws4]:
        for col in range(1, sheet.max_column + 1):
            max_len = 0
            for row in range(1, sheet.max_row + 1):
                val = sheet.cell(row=row, column=col).value
                if val is None:
                    continue
                max_len = max(max_len, len(str(val)))
            sheet.column_dimensions[get_column_letter(col)].width = min(max_len + 2, 50)

    bio = BytesIO()
    wb.save(bio)
    bio.seek(0)
    return bio.getvalue()


def export_report_pptx(markdown_text: str, title: str = "Rapport Qualité Données") -> bytes:
    prs = Presentation()

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = datetime.now().strftime("Généré le %Y-%m-%d %H:%M")

    sections: List[str] = []
    current = []
    for line in markdown_text.splitlines():
        if line.strip().startswith("╔") and current:
            sections.append("\n".join(current).strip())
            current = [line]
        else:
            current.append(line)
    if current:
        sections.append("\n".join(current).strip())

    sections = sections[:4] if sections else [markdown_text[:2000]]

    content_layout = prs.slide_layouts[5]
    for idx, sec in enumerate(sections, 1):
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = f"Section {idx}"

        left = Inches(0.7)
        top = Inches(1.4)
        width = Inches(8.6)
        height = Inches(5.3)
        txBox = s.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        sec_lines = sec.splitlines()[:35]
        tf.text = sec_lines[0] if sec_lines else ""
        tf.paragraphs[0].font.size = Pt(12)

        for ln in sec_lines[1:]:
            p = tf.add_paragraph()
            p.text = ln
            p.font.size = Pt(12)

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()
EOF