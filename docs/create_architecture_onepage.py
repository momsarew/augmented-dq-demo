"""
Schéma d'Architecture ONE-PAGE - Style Blueprint/Poster
Framework Data Quality Probabiliste
Synthèse complète sur une seule page
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
import os

# ============================================================================
# COULEURS
# ============================================================================
BLUE_UI = RGBColor(66, 133, 244)
GREEN_BIZ = RGBColor(52, 168, 83)
PURPLE_DATA = RGBColor(142, 68, 173)
ORANGE_INPUT = RGBColor(251, 188, 4)
RED_SEC = RGBColor(234, 67, 53)
GRAY = RGBColor(95, 99, 104)
DARK_BLUE = RGBColor(26, 35, 126)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
LIGHT_GRAY = RGBColor(245, 245, 250)


def add_box(slide, left, top, width, height, text, fill_color, text_color=WHITE,
            font_size=8, bold=True, border_color=None):
    """Ajoute une boîte"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.color.rgb = fill_color

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    return shape


def add_rect(slide, left, top, width, height, fill_color, border_color=None, opacity=1):
    """Ajoute un rectangle de fond"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()

    return shape


def add_arrow(slide, x, y, direction="down", size=0.2, color=GRAY):
    """Ajoute une flèche"""
    if direction == "down":
        shape_type = MSO_SHAPE.DOWN_ARROW
        w, h = size * 0.7, size
    elif direction == "right":
        shape_type = MSO_SHAPE.RIGHT_ARROW
        w, h = size, size * 0.7
    elif direction == "left":
        shape_type = MSO_SHAPE.LEFT_ARROW
        w, h = size, size * 0.7
    else:
        shape_type = MSO_SHAPE.UP_ARROW
        w, h = size * 0.7, size

    arrow = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


def add_label(slide, left, top, width, height, text, color, font_size=7):
    """Ajoute un label de section"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_connector(slide, x1, y1, x2, y2, color=GRAY):
    """Ajoute une ligne de connexion"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    return connector


def create_onepage_architecture():
    """Crée le schéma d'architecture one-page"""
    prs = Presentation()
    # Format A3 paysage pour plus de détails
    prs.slide_width = Inches(16.5)
    prs.slide_height = Inches(11.7)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = add_rect(slide, 0, 0, 16.5, 11.7, WHITE)

    # =========================================================================
    # TITRE
    # =========================================================================
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(10), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯 Framework Data Quality Probabiliste - Architecture Technique"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    version_box = slide.shapes.add_textbox(Inches(12), Inches(0.2), Inches(4), Inches(0.4))
    tf = version_box.text_frame
    p = tf.paragraphs[0]
    p.text = "v1.2.0 | TOGAF/C4 Model"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SECTION GAUCHE: ACTEURS & ENTRÉES
    # =========================================================================
    add_label(slide, 0.3, 0.8, 2.5, 0.3, "👥 ACTEURS", BLUE_UI)

    add_box(slide, 0.3, 1.2, 1.1, 0.6, "👤\nAnalyst", BLUE_UI, font_size=7)
    add_box(slide, 1.5, 1.2, 1.1, 0.6, "🔒\nDPO", BLUE_UI, font_size=7)
    add_box(slide, 0.3, 1.9, 1.1, 0.6, "💼\nDSI", BLUE_UI, font_size=7)
    add_box(slide, 1.5, 1.9, 1.1, 0.6, "📊\nMétier", BLUE_UI, font_size=7)

    add_label(slide, 0.3, 2.7, 2.5, 0.3, "📥 ENTRÉES", ORANGE_INPUT, font_size=7)
    add_box(slide, 0.3, 3.1, 1.1, 0.5, "📄 CSV", ORANGE_INPUT, BLACK, font_size=8)
    add_box(slide, 1.5, 3.1, 1.1, 0.5, "📊 Excel", ORANGE_INPUT, BLACK, font_size=8)

    # Flèches vers système
    add_arrow(slide, 2.9, 1.5, "right", 0.3, GRAY)
    add_arrow(slide, 2.9, 3.2, "right", 0.3, GRAY)

    # =========================================================================
    # SECTION CENTRALE: SYSTÈME PRINCIPAL
    # =========================================================================
    # Cadre principal du système
    system_frame = add_rect(slide, 3.3, 0.8, 9.5, 7.2, RGBColor(250, 250, 255), DARK_BLUE)

    sys_title = slide.shapes.add_textbox(Inches(3.5), Inches(0.85), Inches(9), Inches(0.3))
    tf = sys_title.text_frame
    p = tf.paragraphs[0]
    p.text = "⚙️ FRAMEWORK DQ PROBABILISTE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # ---- COUCHE UI (Présentation) ----
    add_rect(slide, 3.5, 1.25, 9.1, 1.1, RGBColor(220, 235, 255), BLUE_UI)
    add_label(slide, 3.5, 1.25, 0.5, 1.1, "UI", BLUE_UI, font_size=8)

    ui_components = [
        ("🔍\nScan", 4.1), ("📊\nDashboard", 4.85), ("🎯\nVecteurs", 5.6),
        ("⚠️\nPriorités", 6.35), ("🎚️\nÉlicitation", 7.1), ("🎭\nProfil", 7.85),
        ("🔄\nLineage", 8.6), ("📈\nDAMA", 9.35), ("📋\nReporting", 10.1),
        ("📜\nHistorique", 10.85), ("⚙️\nParams", 11.6)
    ]
    for text, x in ui_components:
        add_box(slide, x, 1.35, 0.7, 0.9, text, BLUE_UI, font_size=6)

    # Flèches UI -> Business
    add_arrow(slide, 8, 2.4, "down", 0.25, GRAY)

    # ---- COUCHE MÉTIER (Business Logic) ----
    add_rect(slide, 3.5, 2.7, 9.1, 1.4, RGBColor(220, 245, 225), GREEN_BIZ)
    add_label(slide, 3.5, 2.7, 0.5, 1.4, "BIZ", GREEN_BIZ, font_size=8)

    add_box(slide, 4.1, 2.85, 1.5, 1.1, "analyzer.py\n─────────\nProfiling\nStatistiques\nTypes", GREEN_BIZ, font_size=6)
    add_box(slide, 5.7, 2.85, 1.5, 1.1, "beta_calculator\n─────────\nVecteurs 4D\nP_DB P_DP\nP_BR P_UP", GREEN_BIZ, font_size=6)
    add_box(slide, 7.3, 2.85, 1.5, 1.1, "ahp_elicitor\n─────────\nPondération\nAHP\nPresets", GREEN_BIZ, font_size=6)
    add_box(slide, 8.9, 2.85, 1.5, 1.1, "risk_scorer\n─────────\nScoring\nPriorisation\nSeuils", GREEN_BIZ, font_size=6)
    add_box(slide, 10.5, 2.85, 1.5, 1.1, "comparator\n─────────\nDAMA vs 4D\nMapping\nBenchmark", GREEN_BIZ, font_size=6)

    # Flèches Business -> Data
    add_arrow(slide, 8, 4.15, "down", 0.25, GRAY)

    # ---- COUCHE DATA ----
    add_rect(slide, 3.5, 4.45, 9.1, 1.2, RGBColor(240, 230, 250), PURPLE_DATA)
    add_label(slide, 3.5, 4.45, 0.5, 1.2, "DATA", PURPLE_DATA, font_size=7)

    add_box(slide, 4.1, 4.6, 2.2, 0.9, "Anomaly Catalog\n─────────────\n60+ règles\nCore + Extended", PURPLE_DATA, font_size=6)
    add_box(slide, 6.4, 4.6, 2.2, 0.9, "Session State\n─────────────\nDataFrame\nRésultats cache", PURPLE_DATA, font_size=6)
    add_box(slide, 8.7, 4.6, 2.2, 0.9, "Audit Trail\n─────────────\nJSON persistant\nTraçabilité", PURPLE_DATA, font_size=6)
    add_box(slide, 11, 4.6, 1.4, 0.9, "Secrets\n────────\nAPI Keys\n🔐", PURPLE_DATA, font_size=6)

    # ---- COUCHE SÉCURITÉ (Transverse) ----
    add_rect(slide, 3.5, 5.85, 9.1, 1, RGBColor(255, 235, 235), RED_SEC)
    add_label(slide, 3.5, 5.85, 0.5, 1, "SEC", RED_SEC, font_size=7)

    add_box(slide, 4.1, 5.95, 1.8, 0.8, "🔒 security.py\nValidation\nSanitization", RED_SEC, font_size=6)
    add_box(slide, 6, 5.95, 1.8, 0.8, "🛡️ XSS Protect\nescape_html()\nPatterns", RED_SEC, font_size=6)
    add_box(slide, 7.9, 5.95, 1.8, 0.8, "📁 File Valid\nMIME | Size\nHash SHA-256", RED_SEC, font_size=6)
    add_box(slide, 9.8, 5.95, 1.5, 0.8, "🔑 Admin\nAuth\nPassword", RED_SEC, font_size=6)
    add_box(slide, 11.4, 5.95, 1, 0.8, "📜\nAudit\nLog", RED_SEC, font_size=6)

    # ---- MODÈLE 4D (encart) ----
    add_rect(slide, 3.5, 7.05, 5.5, 0.85, RGBColor(245, 250, 255), DARK_BLUE)

    model_title = slide.shapes.add_textbox(Inches(3.6), Inches(7.1), Inches(2), Inches(0.2))
    tf = model_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯 MODÈLE 4D"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    add_box(slide, 3.6, 7.35, 1.2, 0.45, "P_DB\nStructure", RGBColor(66, 133, 244), font_size=6)
    add_box(slide, 4.85, 7.35, 1.2, 0.45, "P_DP\nProcess", RGBColor(251, 188, 4), BLACK, font_size=6)
    add_box(slide, 6.1, 7.35, 1.2, 0.45, "P_BR\nRègles", RGBColor(52, 168, 83), font_size=6)
    add_box(slide, 7.35, 7.35, 1.2, 0.45, "P_UP\nUsage", RGBColor(142, 68, 173), font_size=6)

    # Formule
    formula = slide.shapes.add_textbox(Inches(9.2), Inches(7.15), Inches(3.5), Inches(0.6))
    tf = formula.text_frame
    p = tf.paragraphs[0]
    p.text = "Score = Σ(wᵢ×Pᵢ) × mult"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # =========================================================================
    # SECTION DROITE: SERVICES EXTERNES
    # =========================================================================
    add_label(slide, 13, 0.8, 3.2, 0.3, "☁️ SERVICES EXTERNES", GRAY)

    add_box(slide, 13, 1.2, 3.2, 1.1, "🤖 ANTHROPIC\n─────────────\nClaude API\nSonnet 4\nRapports IA", DARK_BLUE, font_size=7)

    add_box(slide, 13, 2.5, 3.2, 1.1, "☁️ STREAMLIT CLOUD\n─────────────\nHébergement\nSSL/HTTPS\nSecrets Manager", GRAY, font_size=7)

    add_box(slide, 13, 3.8, 3.2, 1.1, "📦 GITHUB\n─────────────\nCode source\nCI/CD\nVersioning", RGBColor(70, 70, 80), font_size=7)

    # Flèches système -> externes
    add_arrow(slide, 12.65, 1.6, "right", 0.3, GRAY)
    add_arrow(slide, 12.65, 2.9, "right", 0.3, GRAY)
    add_arrow(slide, 12.65, 4.2, "right", 0.3, GRAY)

    # =========================================================================
    # SECTION BAS: FLUX & SORTIES
    # =========================================================================
    add_label(slide, 0.3, 4, 2.5, 0.3, "📤 SORTIES", GREEN_BIZ, font_size=7)

    add_box(slide, 0.3, 4.4, 1.1, 0.5, "📊 Excel", GREEN_BIZ, font_size=7)
    add_box(slide, 1.5, 4.4, 1.1, 0.5, "📝 MD", GREEN_BIZ, font_size=7)
    add_box(slide, 0.3, 5, 1.1, 0.5, "📄 PDF", GREEN_BIZ, font_size=7)
    add_box(slide, 1.5, 5, 1.1, 0.5, "📋 JSON", GREEN_BIZ, font_size=7)

    # Flèche système -> sorties
    add_arrow(slide, 2.9, 4.6, "left", 0.3, GRAY)

    # =========================================================================
    # FLUX DE DONNÉES (en bas)
    # =========================================================================
    add_rect(slide, 0.3, 8.2, 15.9, 1.8, RGBColor(252, 252, 255), DARK_BLUE)

    flow_title = slide.shapes.add_textbox(Inches(0.5), Inches(8.3), Inches(4), Inches(0.3))
    tf = flow_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🔄 FLUX DE DONNÉES (Pipeline)"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Étapes du flux
    steps = [
        ("1️⃣\nINPUT\nCSV/Excel", ORANGE_INPUT, BLACK),
        ("2️⃣\nVALIDATION\nSécurité", RED_SEC, WHITE),
        ("3️⃣\nANALYSE\nStats", GREEN_BIZ, WHITE),
        ("4️⃣\nDÉTECTION\nAnomalies", PURPLE_DATA, WHITE),
        ("5️⃣\nVECTEURS\n4D", GREEN_BIZ, WHITE),
        ("6️⃣\nPONDÉRATION\nAHP", GREEN_BIZ, WHITE),
        ("7️⃣\nSCORING\nRisques", GREEN_BIZ, WHITE),
        ("8️⃣\nIA\nRapports", DARK_BLUE, WHITE),
        ("9️⃣\nEXPORT\nSorties", BLUE_UI, WHITE),
    ]

    x_start = 0.5
    for i, (text, color, txt_color) in enumerate(steps):
        add_box(slide, x_start + i * 1.7, 8.7, 1.5, 1.1, text, color, txt_color, font_size=7)
        if i < len(steps) - 1:
            add_arrow(slide, x_start + i * 1.7 + 1.55, 9.15, "right", 0.25, GRAY)

    # =========================================================================
    # LÉGENDE
    # =========================================================================
    add_rect(slide, 13, 5.2, 3.2, 2.7, RGBColor(250, 250, 250), GRAY)

    legend_title = slide.shapes.add_textbox(Inches(13.1), Inches(5.3), Inches(3), Inches(0.3))
    tf = legend_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📖 LÉGENDE"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    legend_items = [
        (BLUE_UI, "Couche Présentation (UI)"),
        (GREEN_BIZ, "Couche Métier (Business)"),
        (PURPLE_DATA, "Couche Données (Data)"),
        (RED_SEC, "Couche Sécurité (Transverse)"),
        (ORANGE_INPUT, "Entrées/Données"),
        (GRAY, "Infrastructure"),
        (DARK_BLUE, "Services IA"),
    ]

    y_leg = 5.65
    for color, label in legend_items:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(13.2), Inches(y_leg), Inches(0.25), Inches(0.25))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        txt = slide.shapes.add_textbox(Inches(13.5), Inches(y_leg), Inches(2.5), Inches(0.25))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(7)
        p.font.color.rgb = BLACK

        y_leg += 0.32

    # =========================================================================
    # INFORMATIONS (coin bas droite)
    # =========================================================================
    info_box = slide.shapes.add_textbox(Inches(13), Inches(10.5), Inches(3.2), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Standards: TOGAF / C4 Model / Clean Architecture"
    p.font.size = Pt(7)
    p.font.color.rgb = GRAY
    p = tf.add_paragraph()
    p.text = "github.com/momsarew/augmented-dq-demo"
    p.font.size = Pt(7)
    p.font.color.rgb = GRAY

    # =========================================================================
    # SAVE
    # =========================================================================
    output_path = os.path.join(os.path.dirname(__file__), "Architecture_OnePage.pptx")
    prs.save(output_path)
    print(f"✅ Schéma ONE-PAGE sauvegardé : {output_path}")
    return output_path


if __name__ == "__main__":
    create_onepage_architecture()
