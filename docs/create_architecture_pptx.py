"""
Script de génération du PowerPoint Architecture TOGAF-style
Framework Data Quality Probabiliste
Avec vrais diagrammes visuels et connexions
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import os

# ============================================================================
# COULEURS TOGAF-STYLE
# ============================================================================
# Couche Présentation (bleu)
BLUE_PRESENTATION = RGBColor(66, 133, 244)
BLUE_LIGHT = RGBColor(130, 177, 255)

# Couche Métier (vert)
GREEN_BUSINESS = RGBColor(52, 168, 83)
GREEN_LIGHT = RGBColor(129, 201, 149)

# Couche Application (orange)
ORANGE_APP = RGBColor(251, 188, 4)
ORANGE_LIGHT = RGBColor(255, 214, 102)

# Couche Data (violet)
PURPLE_DATA = RGBColor(142, 68, 173)
PURPLE_LIGHT = RGBColor(187, 143, 206)

# Couche Infrastructure (gris)
GRAY_INFRA = RGBColor(95, 99, 104)
GRAY_LIGHT = RGBColor(154, 160, 166)

# Sécurité (rouge)
RED_SECURITY = RGBColor(234, 67, 53)
RED_LIGHT = RGBColor(244, 143, 137)

# Externe (bleu foncé)
DARK_BLUE = RGBColor(26, 35, 126)

# Fond et texte
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_BG = RGBColor(30, 30, 40)
LIGHT_BG = RGBColor(245, 245, 250)


def add_box(slide, left, top, width, height, text, fill_color, text_color=WHITE,
            font_size=11, bold=True, border_color=None, icon=""):
    """Ajoute une boîte avec texte centré"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.color.rgb = fill_color

    # Texte
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    if icon:
        p.text = f"{icon}\n{text}"
    else:
        p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    return shape


def add_layer_label(slide, left, top, text, color):
    """Ajoute un label de couche vertical"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(0.4), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    return shape


def add_arrow_line(slide, start_x, start_y, end_x, end_y, color=GRAY_INFRA, dashed=False):
    """Ajoute une ligne avec flèche"""
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW if end_x > start_x else MSO_SHAPE.DOWN_ARROW,
        Inches(min(start_x, end_x)),
        Inches(min(start_y, end_y)),
        Inches(abs(end_x - start_x)) if abs(end_x - start_x) > 0.1 else Inches(0.15),
        Inches(abs(end_y - start_y)) if abs(end_y - start_y) > 0.1 else Inches(0.15)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = color
    connector.line.fill.background()
    return connector


def add_connector_line(slide, x1, y1, x2, y2, color=GRAY_INFRA):
    """Ajoute une simple ligne de connexion"""
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(2)
    return line


def add_title(slide, text, subtitle=""):
    """Ajoute titre et sous-titre"""
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(9.4), Inches(0.3))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_INFRA


def add_legend(slide, items, start_x, start_y):
    """Ajoute une légende"""
    y = start_y
    for color, label in items:
        # Carré de couleur
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(start_x), Inches(y), Inches(0.2), Inches(0.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        # Label
        text_box = slide.shapes.add_textbox(
            Inches(start_x + 0.25), Inches(y), Inches(1.5), Inches(0.25)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.color.rgb = BLACK

        y += 0.3


def create_togaf_pptx():
    """Crée la présentation TOGAF-style"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # SLIDE 1: TITRE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    # Titre principal
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Framework Data Quality Probabiliste"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Sous-titre
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Architecture Technique - Style TOGAF/ArchiMate"
    p.font.size = Pt(24)
    p.font.color.rgb = BLUE_PRESENTATION
    p.alignment = PP_ALIGN.CENTER

    # Version
    ver_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11), Inches(0.4))
    tf = ver_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Version 1.2.0 | Février 2025"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY_INFRA
    p.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 2: VUE CONTEXTE (C4 Level 1)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    add_title(slide, "Vue Contexte Système", "C4 Model - Level 1")

    # Acteurs (à gauche)
    add_box(slide, 0.5, 1.5, 2, 1, "Data Analyst", BLUE_PRESENTATION, icon="👤")
    add_box(slide, 0.5, 2.8, 2, 1, "DPO / DSI", BLUE_PRESENTATION, icon="👥")
    add_box(slide, 0.5, 4.1, 2, 1, "Direction Métier", BLUE_PRESENTATION, icon="👔")

    # Système central
    add_box(slide, 4.5, 2.5, 4, 2.5, "Framework DQ\nProbabiliste\n\n🎯 Analyse qualité\n📊 Vecteurs 4D\n📋 Rapports IA",
            GREEN_BUSINESS, font_size=14, border_color=RGBColor(30, 100, 50))

    # Systèmes externes (à droite)
    add_box(slide, 10, 1.5, 2.5, 1, "Claude API\n(Anthropic)", PURPLE_DATA, icon="🤖")
    add_box(slide, 10, 3, 2.5, 1, "Streamlit Cloud", GRAY_INFRA, icon="☁️")
    add_box(slide, 10, 4.5, 2.5, 1, "GitHub", GRAY_INFRA, icon="📦")

    # Données (en bas)
    add_box(slide, 4.5, 5.8, 1.8, 0.9, "CSV", ORANGE_APP, icon="📄")
    add_box(slide, 6.7, 5.8, 1.8, 0.9, "Excel", ORANGE_APP, icon="📊")

    # Flèches (représentées par des petites formes)
    # Utilisateurs -> Système
    add_connector_line(slide, 2.5, 2, 4.5, 3)
    add_connector_line(slide, 2.5, 3.3, 4.5, 3.5)
    add_connector_line(slide, 2.5, 4.6, 4.5, 4)

    # Système -> Externes
    add_connector_line(slide, 8.5, 3, 10, 2)
    add_connector_line(slide, 8.5, 3.5, 10, 3.5)
    add_connector_line(slide, 8.5, 4, 10, 5)

    # Données -> Système
    add_connector_line(slide, 5.4, 5.8, 5.4, 5)
    add_connector_line(slide, 7.6, 5.8, 7.6, 5)

    # Légende
    add_legend(slide, [
        (BLUE_PRESENTATION, "Acteurs"),
        (GREEN_BUSINESS, "Système"),
        (PURPLE_DATA, "Service IA"),
        (GRAY_INFRA, "Infrastructure"),
        (ORANGE_APP, "Données"),
    ], 0.5, 5.8)

    # =========================================================================
    # SLIDE 3: ARCHITECTURE EN COUCHES (TOGAF)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    add_title(slide, "Architecture en Couches", "TOGAF - Vue Applicative")

    # ---- COUCHE PRÉSENTATION ----
    add_layer_label(slide, 0.3, 1.2, "UI", BLUE_PRESENTATION)

    # Barre de couche
    layer1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(12), Inches(1.3))
    layer1.fill.solid()
    layer1.fill.fore_color.rgb = BLUE_LIGHT
    layer1.line.color.rgb = BLUE_PRESENTATION
    layer1.line.width = Pt(2)

    # Composants UI
    add_box(slide, 1, 1.35, 1.5, 1, "🔍 Scan", BLUE_PRESENTATION, font_size=10)
    add_box(slide, 2.6, 1.35, 1.5, 1, "📊 Dashboard", BLUE_PRESENTATION, font_size=10)
    add_box(slide, 4.2, 1.35, 1.5, 1, "🎯 Vecteurs", BLUE_PRESENTATION, font_size=10)
    add_box(slide, 5.8, 1.35, 1.5, 1, "⚠️ Priorités", BLUE_PRESENTATION, font_size=10)
    add_box(slide, 7.4, 1.35, 1.5, 1, "📈 DAMA", BLUE_PRESENTATION, font_size=10)
    add_box(slide, 9, 1.35, 1.5, 1, "📋 Reporting", BLUE_PRESENTATION, font_size=10)
    add_box(slide, 10.6, 1.35, 1.5, 1, "📜 Historique", BLUE_PRESENTATION, font_size=10)

    # ---- COUCHE MÉTIER ----
    add_layer_label(slide, 0.3, 2.7, "BIZ", GREEN_BUSINESS)

    layer2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.7), Inches(12), Inches(1.5))
    layer2.fill.solid()
    layer2.fill.fore_color.rgb = GREEN_LIGHT
    layer2.line.color.rgb = GREEN_BUSINESS
    layer2.line.width = Pt(2)

    # Composants Métier
    add_box(slide, 1, 2.9, 2.2, 1.1, "analyzer.py\n─────────\nProfiling\nStatistiques", GREEN_BUSINESS, font_size=9)
    add_box(slide, 3.4, 2.9, 2.2, 1.1, "beta_calculator.py\n─────────\nVecteurs 4D\nP_DB, P_DP, P_BR, P_UP", GREEN_BUSINESS, font_size=9)
    add_box(slide, 5.8, 2.9, 2.2, 1.1, "risk_scorer.py\n─────────\nScoring\nPriorisation", GREEN_BUSINESS, font_size=9)
    add_box(slide, 8.2, 2.9, 2.2, 1.1, "ahp_elicitor.py\n─────────\nPondération\nAHP", GREEN_BUSINESS, font_size=9)
    add_box(slide, 10.6, 2.9, 2.2, 1.1, "comparator.py\n─────────\nDAMA vs 4D\nBenchmark", GREEN_BUSINESS, font_size=9)

    # ---- COUCHE DATA ----
    add_layer_label(slide, 0.3, 4.4, "DATA", PURPLE_DATA)

    layer3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.4), Inches(12), Inches(1.3))
    layer3.fill.solid()
    layer3.fill.fore_color.rgb = PURPLE_LIGHT
    layer3.line.color.rgb = PURPLE_DATA
    layer3.line.width = Pt(2)

    add_box(slide, 1, 4.55, 3, 1, "Anomaly Catalog\n─────────\n60+ règles détection", PURPLE_DATA, font_size=10)
    add_box(slide, 4.2, 4.55, 3, 1, "Session State\n─────────\nDataFrame, Résultats", PURPLE_DATA, font_size=10)
    add_box(slide, 7.4, 4.55, 3, 1, "Audit Trail\n─────────\nJSON persistant", PURPLE_DATA, font_size=10)
    add_box(slide, 10.6, 4.55, 2.2, 1, "Secrets\n─────────\nAPI Keys", PURPLE_DATA, font_size=10)

    # ---- COUCHE TRANSVERSE (Sécurité) ----
    sec_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(5.9), Inches(12.5), Inches(1))
    sec_box.fill.solid()
    sec_box.fill.fore_color.rgb = RED_LIGHT
    sec_box.line.color.rgb = RED_SECURITY
    sec_box.line.width = Pt(2)

    add_box(slide, 0.5, 6.05, 2.5, 0.7, "🔐 security.py\nValidation | Sanitization", RED_SECURITY, font_size=9)
    add_box(slide, 3.2, 6.05, 2.5, 0.7, "🔒 XSS Protection\nescape_html()", RED_SECURITY, font_size=9)
    add_box(slide, 5.9, 6.05, 2.5, 0.7, "📁 File Validation\nMIME | Size | Hash", RED_SECURITY, font_size=9)
    add_box(slide, 8.6, 6.05, 2, 0.7, "🔑 Auth Admin\nPassword", RED_SECURITY, font_size=9)
    add_box(slide, 10.8, 6.05, 2, 0.7, "📜 Audit Trail\nTraçabilité", RED_SECURITY, font_size=9)

    # =========================================================================
    # SLIDE 4: FLUX DE DONNÉES
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    add_title(slide, "Flux de Données", "Pipeline de traitement Data Quality")

    # Étape 1: Input
    add_box(slide, 0.5, 1.5, 1.8, 1.2, "📥 INPUT\n─────\nCSV / Excel", ORANGE_APP, BLACK, font_size=10)

    # Étape 2: Validation
    add_box(slide, 2.8, 1.5, 1.8, 1.2, "🔒 VALIDATION\n─────\nTaille\nMIME\nHash", RED_SECURITY, font_size=10)

    # Étape 3: Analyse
    add_box(slide, 5.1, 1.5, 1.8, 1.2, "🔍 ANALYSE\n─────\nStats\nProfiling", GREEN_BUSINESS, font_size=10)

    # Étape 4: Détection
    add_box(slide, 7.4, 1.5, 1.8, 1.2, "⚠️ DÉTECTION\n─────\n60+ règles\nAnomalies", PURPLE_DATA, font_size=10)

    # Étape 5: Vecteurs
    add_box(slide, 9.7, 1.5, 1.8, 1.2, "🎯 VECTEURS\n─────\nP_DB P_DP\nP_BR P_UP", GREEN_BUSINESS, font_size=10)

    # Étape 6: Scoring
    add_box(slide, 0.5, 3.5, 1.8, 1.2, "📊 SCORING\n─────\nPondération\nAHP", GREEN_BUSINESS, font_size=10)

    # Étape 7: IA
    add_box(slide, 2.8, 3.5, 1.8, 1.2, "🤖 IA\n─────\nClaude API\nRapports", DARK_BLUE, font_size=10)

    # Étape 8: Output
    add_box(slide, 5.1, 3.5, 1.8, 1.2, "📤 OUTPUT\n─────\nDashboard\nExports", BLUE_PRESENTATION, font_size=10)

    # Flèches horizontales (ligne 1)
    for x in [2.3, 4.6, 6.9, 9.2]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(1.95), Inches(0.4), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY_INFRA
        arrow.line.fill.background()

    # Flèche descendante
    arrow_down = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.4), Inches(2.8), Inches(0.3), Inches(0.5))
    arrow_down.fill.solid()
    arrow_down.fill.fore_color.rgb = GRAY_INFRA
    arrow_down.line.fill.background()

    # Flèche retour
    arrow_back = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(9.7), Inches(3.9), Inches(6.9), Inches(0.3))
    arrow_back.fill.solid()
    arrow_back.fill.fore_color.rgb = RGBColor(200, 200, 200)
    arrow_back.line.fill.background()

    # Flèches ligne 2
    for x in [2.3, 4.6]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(3.95), Inches(0.4), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY_INFRA
        arrow.line.fill.background()

    # Audit Trail (en bas)
    audit_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.2), Inches(12), Inches(0.8))
    audit_box.fill.solid()
    audit_box.fill.fore_color.rgb = RGBColor(255, 235, 235)
    audit_box.line.color.rgb = RED_SECURITY
    audit_box.line.width = Pt(1)

    add_box(slide, 5, 5.3, 3, 0.6, "📜 AUDIT TRAIL - Traçabilité complète de chaque étape", RED_SECURITY, WHITE, font_size=11)

    # =========================================================================
    # SLIDE 5: MODÈLE PROBABILISTE 4D
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    add_title(slide, "Modèle Probabiliste 4D", "Vecteur de risque multi-dimensionnel")

    # Les 4 dimensions
    add_box(slide, 1, 1.5, 2.5, 2, "P_DB\n━━━━━━━\n📦 Database\nStructure\n─────\n• Nulls\n• Types\n• Clés",
            RGBColor(66, 133, 244), font_size=11)

    add_box(slide, 4, 1.5, 2.5, 2, "P_DP\n━━━━━━━\n⚙️ Data\nProcessing\n─────\n• Formats\n• Conversions\n• Outliers",
            RGBColor(251, 188, 4), font_size=11)

    add_box(slide, 7, 1.5, 2.5, 2, "P_BR\n━━━━━━━\n📋 Business\nRules\n─────\n• Patterns\n• Plages\n• Logique",
            RGBColor(52, 168, 83), font_size=11)

    add_box(slide, 10, 1.5, 2.5, 2, "P_UP\n━━━━━━━\n🎯 Usage-fit\nPurpose\n─────\n• Contexte\n• Criticité\n• Adéquation",
            RGBColor(142, 68, 173), font_size=11)

    # Flèches vers formule
    for x in [2.25, 5.25, 8.25, 11.25]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(3.6), Inches(0.25), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY_INFRA
        arrow.line.fill.background()

    # Formule centrale
    formula_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(4.2), Inches(9), Inches(1))
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = DARK_BLUE
    formula_box.line.fill.background()

    formula_text = slide.shapes.add_textbox(Inches(2), Inches(4.35), Inches(9), Inches(0.7))
    tf = formula_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Score = Σ (wᵢ × Pᵢ) × multiplicateur_profil"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Profils de risque
    add_box(slide, 1.5, 5.5, 3, 1.2, "🛡️ PRUDENT\n×1.3\nAmplifie les risques", RGBColor(234, 67, 53), font_size=11)
    add_box(slide, 5, 5.5, 3, 1.2, "⚖️ ÉQUILIBRÉ\n×1.0\nNeutre", RGBColor(251, 188, 4), font_size=11)
    add_box(slide, 8.5, 5.5, 3, 1.2, "🚀 AUDACIEUX\n×0.7\nAtténue les risques", RGBColor(52, 168, 83), font_size=11)

    # =========================================================================
    # SLIDE 6: MAPPING DAMA vs 4D
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    add_title(slide, "Mapping DAMA ↔ Probabiliste 4D", "Correspondance des dimensions qualité")

    # Colonne DAMA
    add_box(slide, 0.8, 1.3, 2.5, 0.7, "📐 DAMA", DARK_BLUE, font_size=14)
    add_box(slide, 0.8, 2.1, 2.5, 0.7, "Complétude", BLUE_PRESENTATION, font_size=12)
    add_box(slide, 0.8, 2.9, 2.5, 0.7, "Unicité", BLUE_PRESENTATION, font_size=12)
    add_box(slide, 0.8, 3.7, 2.5, 0.7, "Validité", GREEN_BUSINESS, font_size=12)
    add_box(slide, 0.8, 4.5, 2.5, 0.7, "Cohérence", GREEN_BUSINESS, font_size=12)
    add_box(slide, 0.8, 5.3, 2.5, 0.7, "Fraîcheur", ORANGE_APP, BLACK, font_size=12)
    add_box(slide, 0.8, 6.1, 2.5, 0.7, "Exactitude", PURPLE_DATA, font_size=12)

    # Flèches
    for y in [2.35, 2.35, 3.95, 3.95, 5.55, 6.35]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.4), Inches(y), Inches(0.8), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY_INFRA
        arrow.line.fill.background()

    # Colonne 4D
    add_box(slide, 4.5, 1.3, 2.5, 0.7, "🎯 4D", DARK_BLUE, font_size=14)
    add_box(slide, 4.5, 2.1, 2.5, 1.5, "P_DB\n📦 Structure\nBase de données", RGBColor(66, 133, 244), font_size=11)
    add_box(slide, 4.5, 3.7, 2.5, 1.5, "P_BR\n📋 Règles Métier\nBusiness Rules", RGBColor(52, 168, 83), font_size=11)
    add_box(slide, 4.5, 5.3, 2.5, 0.7, "P_DP\n⚙️ Traitements", RGBColor(251, 188, 4), BLACK, font_size=11)
    add_box(slide, 4.5, 6.1, 2.5, 0.7, "P_UP\n🎯 Usage", RGBColor(142, 68, 173), font_size=11)

    # Avantages 4D
    add_box(slide, 8, 1.5, 4.5, 2, "✅ Avantages 4D\n━━━━━━━━━━━━━━\n• Pondération contextuelle\n• Adapté par usage\n• Profils de risque\n• Multi-dimensionnel",
            GREEN_BUSINESS, font_size=12)

    add_box(slide, 8, 4, 4.5, 2.5, "📊 Usages supportés\n━━━━━━━━━━━━━━\n• Paie / Réglementaire\n• Reporting Social\n• Dashboard Opérationnel\n• Analyse RH\n• Custom...",
            BLUE_PRESENTATION, font_size=12)

    # =========================================================================
    # SLIDE 7: ARCHITECTURE DÉPLOIEMENT
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    add_title(slide, "Architecture de Déploiement", "Infrastructure Cloud")

    # Utilisateur
    add_box(slide, 0.5, 2.5, 2, 1.5, "👤 Utilisateur\n─────────\nNavigateur Web\nHTTPS", BLUE_PRESENTATION, font_size=11)

    # Flèche
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.6), Inches(3.1), Inches(0.8), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY_INFRA
    arrow.line.fill.background()

    # Streamlit Cloud
    cloud_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.8), Inches(1.3), Inches(5.5), Inches(5))
    cloud_box.fill.solid()
    cloud_box.fill.fore_color.rgb = RGBColor(240, 240, 250)
    cloud_box.line.color.rgb = GRAY_INFRA
    cloud_box.line.width = Pt(2)

    cloud_label = slide.shapes.add_textbox(Inches(4), Inches(1.4), Inches(5), Inches(0.4))
    tf = cloud_label.text_frame
    p = tf.paragraphs[0]
    p.text = "☁️ STREAMLIT CLOUD"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GRAY_INFRA

    add_box(slide, 4, 2, 2.3, 1, "Load Balancer\n─────────\nSSL/HTTPS", GRAY_INFRA, font_size=10)
    add_box(slide, 6.8, 2, 2.3, 1, "Container\n─────────\nPython 3.11", GRAY_INFRA, font_size=10)
    add_box(slide, 4, 3.3, 2.3, 1.3, "app.py\n─────────\nStreamlit Server\n12 onglets", GREEN_BUSINESS, font_size=10)
    add_box(slide, 6.8, 3.3, 2.3, 1.3, "backend/\n─────────\nEngine\nSecurity\nAudit", GREEN_BUSINESS, font_size=10)
    add_box(slide, 4, 4.9, 2.3, 1, "Session State\n─────────\nRAM", PURPLE_DATA, font_size=10)
    add_box(slide, 6.8, 4.9, 2.3, 1, "Secrets\n─────────\nAPI Keys 🔐", RED_SECURITY, font_size=10)

    # Services externes
    add_box(slide, 10.5, 2, 2.3, 1.2, "🤖 Anthropic\n─────────\nClaude API\nSonnet 4", DARK_BLUE, font_size=10)
    add_box(slide, 10.5, 3.8, 2.3, 1.2, "📦 GitHub\n─────────\nCode Source\nCI/CD", GRAY_INFRA, font_size=10)

    # Flèches externes
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.4), Inches(2.4), Inches(0.9), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY_INFRA
    arrow.line.fill.background()

    arrow = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(9.4), Inches(4.2), Inches(0.9), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY_INFRA
    arrow.line.fill.background()

    # =========================================================================
    # SLIDE 8: STACK TECHNIQUE
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    add_title(slide, "Stack Technique", "Technologies utilisées")

    # Frontend
    add_box(slide, 0.5, 1.3, 2, 0.5, "FRONTEND", BLUE_PRESENTATION, font_size=12)
    add_box(slide, 0.5, 1.9, 2, 0.6, "Streamlit\n1.53.1", BLUE_LIGHT, BLACK, font_size=10, bold=False)
    add_box(slide, 0.5, 2.6, 2, 0.6, "Plotly\n6.5.2", BLUE_LIGHT, BLACK, font_size=10, bold=False)

    # Data
    add_box(slide, 2.8, 1.3, 2, 0.5, "DATA", PURPLE_DATA, font_size=12)
    add_box(slide, 2.8, 1.9, 2, 0.6, "Pandas\n2.3.3", PURPLE_LIGHT, BLACK, font_size=10, bold=False)
    add_box(slide, 2.8, 2.6, 2, 0.6, "NumPy\n2.4.1", PURPLE_LIGHT, BLACK, font_size=10, bold=False)
    add_box(slide, 2.8, 3.3, 2, 0.6, "SciPy\n1.17.0", PURPLE_LIGHT, BLACK, font_size=10, bold=False)

    # Export
    add_box(slide, 5.1, 1.3, 2, 0.5, "EXPORT", ORANGE_APP, BLACK, font_size=12)
    add_box(slide, 5.1, 1.9, 2, 0.6, "OpenPyXL\n3.1.5", ORANGE_LIGHT, BLACK, font_size=10, bold=False)
    add_box(slide, 5.1, 2.6, 2, 0.6, "ReportLab\n4.4.9", ORANGE_LIGHT, BLACK, font_size=10, bold=False)

    # IA
    add_box(slide, 7.4, 1.3, 2, 0.5, "IA", DARK_BLUE, font_size=12)
    add_box(slide, 7.4, 1.9, 2, 0.6, "Anthropic\n0.76.0", RGBColor(100, 110, 180), font_size=10, bold=False)
    add_box(slide, 7.4, 2.6, 2, 0.6, "Claude\nSonnet 4", RGBColor(100, 110, 180), font_size=10, bold=False)

    # Infra
    add_box(slide, 9.7, 1.3, 2, 0.5, "INFRA", GRAY_INFRA, font_size=12)
    add_box(slide, 9.7, 1.9, 2, 0.6, "Streamlit\nCloud", GRAY_LIGHT, BLACK, font_size=10, bold=False)
    add_box(slide, 9.7, 2.6, 2, 0.6, "GitHub", GRAY_LIGHT, BLACK, font_size=10, bold=False)

    # Principes architecturaux
    add_box(slide, 0.5, 4.3, 5.8, 2.5,
            "🏗️ Principes Architecturaux\n━━━━━━━━━━━━━━━━━━━━━━━━\n\n"
            "• Séparation des responsabilités (SoC)\n"
            "• Couches indépendantes\n"
            "• Sécurité by design\n"
            "• Clean Architecture\n"
            "• Extensibilité (catalogs, profils)",
            RGBColor(240, 245, 250), DARK_BLUE, font_size=11, border_color=DARK_BLUE)

    add_box(slide, 6.8, 4.3, 5.8, 2.5,
            "📊 Métriques\n━━━━━━━━━━━━━━━━━━━━━━━━\n\n"
            "• 60+ règles d'anomalies\n"
            "• 12 onglets fonctionnels\n"
            "• 4 dimensions probabilistes\n"
            "• 6 dimensions DAMA mappées\n"
            "• 100% traçabilité (Audit Trail)",
            RGBColor(240, 250, 245), DARK_BLUE, font_size=11, border_color=GREEN_BUSINESS)

    # =========================================================================
    # SAVE
    # =========================================================================
    output_path = os.path.join(os.path.dirname(__file__), "Architecture_Framework_DQ_TOGAF.pptx")
    prs.save(output_path)
    print(f"✅ Présentation TOGAF sauvegardée : {output_path}")
    return output_path


if __name__ == "__main__":
    create_togaf_pptx()
